"""
Modul untuk generate laporan PowerPoint (.pptx) dari data, chart, dan branding user.
Mendukung 3 gaya template: minimalist, corporate, colorful.
"""
import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from modules.translations import get_labels

WHITE = RGBColor(255, 255, 255)
DARK = RGBColor(30, 30, 30)
GREY = RGBColor(90, 90, 90)
LIGHT_BG = RGBColor(245, 246, 250)

TEMPLATES = {
    "minimalist": {
        "cover_bg": "white",       # cover pakai background putih + aksen garis warna brand
        "heading_color": "brand",
        "card_style": "outline",
    },
    "corporate": {
        "cover_bg": "brand",       # cover full warna brand solid
        "heading_color": "brand",
        "card_style": "filled",
    },
    "colorful": {
        "cover_bg": "brand",
        "heading_color": "brand",
        "card_style": "accent_bar",
    },
}


def _hex_to_rgb(hex_color: str) -> RGBColor:
    hex_color = hex_color.lstrip("#")
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


def _add_cover_slide(prs, title, period, brand_color, logo_bytes, template: str):
    cfg = TEMPLATES.get(template, TEMPLATES["corporate"])
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    brand_rgb = _hex_to_rgb(brand_color)

    if cfg["cover_bg"] == "brand":
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = brand_rgb
        bg.line.fill.background()
        bg.shadow.inherit = False
        title_color = WHITE
        period_color = WHITE
    else:
        # minimalist: background putih + garis aksen di kiri
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.25), prs.slide_height)
        accent.fill.solid()
        accent.fill.fore_color.rgb = brand_rgb
        accent.line.fill.background()
        accent.shadow.inherit = False
        title_color = DARK
        period_color = GREY

    if logo_bytes:
        slide.shapes.add_picture(io.BytesIO(logo_bytes), Inches(0.6), Inches(0.4), height=Inches(0.8))

    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(2.7), prs.slide_width - Inches(1.4), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = title_color

    period_box = slide.shapes.add_textbox(Inches(0.7), Inches(3.8), prs.slide_width - Inches(1.4), Inches(0.6))
    p2 = period_box.text_frame.paragraphs[0]
    p2.text = period
    p2.font.size = Pt(18)
    p2.font.color.rgb = period_color
    return slide


def _add_heading(slide, prs, text, brand_color):
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), prs.slide_width - Inches(1), Inches(0.7))
    p = title_box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = _hex_to_rgb(brand_color)


def _add_summary_slide(prs, summary: dict, brand_color: str, template: str, labels: dict):
    cfg = TEMPLATES.get(template, TEMPLATES["corporate"])
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_heading(slide, prs, labels["summary_heading"], brand_color)

    metrics = list(summary.items())
    n = len(metrics) or 1
    box_width = (prs.slide_width - Inches(1) - Inches(0.3) * (n - 1)) / n
    x = Inches(0.5)
    brand_rgb = _hex_to_rgb(brand_color)

    for name, stats in metrics:
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(1.5), box_width, Inches(2.2))
        card.shadow.inherit = False

        if cfg["card_style"] == "filled":
            card.fill.solid()
            card.fill.fore_color.rgb = brand_rgb
            name_color, val_color, sub_color = WHITE, WHITE, RGBColor(230, 230, 230)
            card.line.fill.background()
        elif cfg["card_style"] == "accent_bar":
            card.fill.solid()
            card.fill.fore_color.rgb = LIGHT_BG
            card.line.color.rgb = brand_rgb
            card.line.width = Pt(2.5)
            name_color, val_color, sub_color = DARK, brand_rgb, GREY
        else:  # outline
            card.fill.solid()
            card.fill.fore_color.rgb = WHITE
            card.line.color.rgb = brand_rgb
            card.line.width = Pt(1)
            name_color, val_color, sub_color = DARK, brand_rgb, GREY

        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.15)
        tf.margin_top = Inches(0.15)

        p1 = tf.paragraphs[0]
        p1.text = name
        p1.font.size = Pt(14)
        p1.font.bold = True
        p1.font.color.rgb = name_color

        p2 = tf.add_paragraph()
        p2.text = f"{labels['total']}: {stats['total']:,.0f}"
        p2.font.size = Pt(16)
        p2.font.color.rgb = val_color
        p2.font.bold = True

        p3 = tf.add_paragraph()
        p3.text = f"{labels['average']}: {stats['average']:,.1f}"
        p3.font.size = Pt(12)
        p3.font.color.rgb = sub_color

        x += box_width + Inches(0.3)
    return slide


def _add_insight_slide(prs, insights: list, brand_color: str, labels: dict):
    if not insights:
        return None
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_heading(slide, prs, labels["insight_heading"], brand_color)

    body = slide.shapes.add_textbox(Inches(0.6), Inches(1.3), prs.slide_width - Inches(1.2), Inches(5.5))
    tf = body.text_frame
    tf.word_wrap = True
    brand_rgb = _hex_to_rgb(brand_color)

    for i, text in enumerate(insights[:6]):
        clean_text = text.replace("**", "")
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"•  {clean_text}"
        p.font.size = Pt(15)
        p.font.color.rgb = DARK
        p.space_after = Pt(12)
    return slide


def _add_chart_slide(prs, title, chart_png_bytes, brand_color):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_heading(slide, prs, title, brand_color)
    img_stream = io.BytesIO(chart_png_bytes)
    slide.shapes.add_picture(img_stream, Inches(0.6), Inches(1.2), width=prs.slide_width - Inches(1.2))
    return slide


def _add_table_slide(prs, df, brand_color, max_rows=12, labels=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_heading(slide, prs, labels["detail_heading"], brand_color)

    display_df = df.head(max_rows)
    rows, cols = display_df.shape[0] + 1, display_df.shape[1]
    table_shape = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.1),
                                          prs.slide_width - Inches(1), Inches(0.35) * rows)
    table = table_shape.table

    for j, col_name in enumerate(display_df.columns):
        cell = table.cell(0, j)
        cell.text = str(col_name)
        cell.fill.solid()
        cell.fill.fore_color.rgb = _hex_to_rgb(brand_color)
        for para in cell.text_frame.paragraphs:
            para.font.size = Pt(11)
            para.font.bold = True
            para.font.color.rgb = WHITE

    for i in range(display_df.shape[0]):
        for j in range(cols):
            cell = table.cell(i + 1, j)
            cell.text = str(display_df.iloc[i, j])
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(10)
    return slide


def _add_watermark(prs):
    """Tambah teks watermark kecil di pojok kanan bawah tiap slide (versi gratis)."""
    for slide in prs.slides:
        box = slide.shapes.add_textbox(
            prs.slide_width - Inches(3.2), prs.slide_height - Inches(0.4), Inches(3), Inches(0.3)
        )
        p = box.text_frame.paragraphs[0]
        p.text = "Dibuat dengan Generator Laporan Otomatis"
        p.font.size = Pt(8)
        p.font.color.rgb = RGBColor(180, 180, 180)
        p.alignment = PP_ALIGN.RIGHT


def generate_pptx(title: str, period: str, brand_color: str, logo_bytes,
                   summary: dict, chart_images: list, df_preview,
                   insights: list = None, template: str = "corporate",
                   lang: str = "id", watermark: bool = False) -> bytes:
    """
    chart_images: list of tuples (chart_title, png_bytes)
    template: "minimalist" | "corporate" | "colorful"
    Return: bytes dari file .pptx
    """
    labels = get_labels(lang)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    _add_cover_slide(prs, title, period, brand_color, logo_bytes, template)
    if summary:
        _add_summary_slide(prs, summary, brand_color, template, labels)
    if insights:
        _add_insight_slide(prs, insights, brand_color, labels)
    for chart_title, png_bytes in chart_images:
        _add_chart_slide(prs, chart_title, png_bytes, brand_color)
    if df_preview is not None and not df_preview.empty:
        _add_table_slide(prs, df_preview, brand_color, labels=labels)

    if watermark:
        _add_watermark(prs)

    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return out.getvalue()
